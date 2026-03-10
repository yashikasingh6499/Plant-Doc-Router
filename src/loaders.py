import io
import os
from typing import Dict, List

import fitz
from docx import Document
from pptx import Presentation

from src.config import ENABLE_VISION_SUMMARY, VISION_PAGE_TEXT_THRESHOLD
from src.utils import clean_text
from src.vision import summarize_page_image


def load_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return clean_text(f.read())


def load_docx(path: str) -> str:
    doc = Document(path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return clean_text("\n\n".join(paragraphs))


def load_pptx(path: str) -> str:
    prs = Presentation(path)
    slides_text = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        parts = [f"Slide {slide_idx}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                txt = shape.text.strip()
                if txt:
                    parts.append(txt)
        slides_text.append("\n".join(parts))

    return clean_text("\n\n".join(slides_text))


def render_pdf_page_to_png(page, zoom: float = 1.5) -> bytes:
    matrix = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=matrix, alpha=False)
    return pix.tobytes("png")


def load_pdf(path: str) -> str:
    doc = fitz.open(path)
    page_blocks = []

    for page_num, page in enumerate(doc, start=1):
        page_text = clean_text(page.get_text("text"))
        page_images = page.get_images(full=True)
        vision_summary = ""

        should_run_vision = ENABLE_VISION_SUMMARY and (
            len(page_images) > 0 or len(page_text) < VISION_PAGE_TEXT_THRESHOLD
        )

        if should_run_vision:
            try:
                image_bytes = render_pdf_page_to_png(page)
                vision_summary = summarize_page_image(image_bytes)
            except Exception:
                vision_summary = ""

        parts = [f"Page {page_num}"]
        if page_text:
            parts.append(page_text)
        if vision_summary:
            parts.append(f"Visual Summary: {vision_summary}")

        page_blocks.append("\n\n".join(parts))

    return clean_text("\n\n".join(page_blocks))


def load_document(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()

    if ext in [".txt", ".md"]:
        return load_txt(path)
    if ext == ".pdf":
        return load_pdf(path)
    if ext == ".docx":
        return load_docx(path)
    if ext == ".pptx":
        return load_pptx(path)

    return ""


def load_source_document(file_path: str) -> List[Dict]:
    if not os.path.exists(file_path):
        return []

    text = load_document(file_path)
    if not text:
        return []

    return [{
        "file_name": os.path.basename(file_path),
        "file_path": file_path,
        "text": text
    }]